import os
import docx
from pptx import Presentation
import PyPDF2
from pdf2image import convert_from_path
from ocr import OCRProcessor  # Import OCRProcessor for collaboration
import pytesseract

class FileHandler:
    
    def __init__(self,tesseract_path="C:/Program Files/Tesseract-OCR/tesseract.exe", poppler_path="C:/poppler-24.08.0/Library/bin"):
        self.ocr_processor = OCRProcessor(tesseract_path,poppler_path) 

    def extract_from_txt(self, file_path):
        """Extract text from TXT files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            return f"Error reading TXT file: {str(e)}"

    def extract_from_docx(self, file_path):
        """Extract text from DOCX files"""
        try:
            doc = docx.Document(file_path)
            full_text = [para.text for para in doc.paragraphs]
            return '\n'.join(full_text)
        except Exception as e:
            return f"Error reading DOCX file: {str(e)}"

    def has_images(self,pdf_path):
      # Open the PDF
        pdf = PyPDF2.PdfReader(pdf_path)
        for page_num in range(len(pdf.pages)):
            page = pdf.pages[page_num]
            # Check for /XObject resources (images are often stored here)
            if '/XObject' in page['/Resources']:
                xobjects = page['/Resources']['/XObject'].get_object()
                for obj in xobjects:
                    if xobjects[obj]['/Subtype'] == '/Image':
                        return True
        # Fallback: Try converting to images (checks rendered content)
        images = convert_from_path(pdf_path)
        return len(images) > 0
        
        
    
    def extract_from_pdf(self, file_path):
        """Extract text from text-based PDFs, redirect to OCR if needed"""
        try:
            reader = PyPDF2.PdfReader(file_path)
            full_text = [page.extract_text() for page in reader.pages]
            extracted_text = '\n'.join(full_text)

            # Check if the extracted text is empty or minimal (likely image-based)
            # if not extracted_text.strip() or len(extracted_text.strip()) < 10:
            if self.has_images(file_path):
                print(f"No text extracted from {file_path}. Redirecting to OCR...")
                return self.ocr_processor.process_pdf(file_path)
            return extracted_text
        except Exception as e:
            return f"Error reading PDF file: {str(e)}"

    def extract_from_ppt(self, file_path):
        """Extract text from PPT/PPTX files"""
        try:
            prs = Presentation(file_path)
            full_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text.append(shape.text)
            return '\n'.join(full_text)
        except Exception as e:
            return f"Error reading PPT file: {str(e)}"

    def process_file(self, file_path):
        """Process any supported file type"""
        if not os.path.exists(file_path):
            return "File not found!"
        
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension == '.txt':
            return self.extract_from_txt(file_path)
        elif file_extension == '.docx':
            return self.extract_from_docx(file_path)
        elif file_extension == '.pdf':
            return self.extract_from_pdf(file_path)
        elif file_extension in ['.ppt', '.pptx']:
            return self.extract_from_ppt(file_path)
        else:
            return "Unsupported file format!"

def main():
    handler = FileHandler(poppler_path="C:/poppler-24.08.0/Library/bin") 
    # Add poppler_path="C:/poppler-23.05.0/bin" if needed
    
    test_files = [
       "dataset.txt",
       "Introduction-to-Probability.pptx",
       "Doc-Comparison SRS.docx",
       "Machine learning algorithms are largely driven by mathematics.pdf"
      
    ]
    
    for file_path in test_files:
        print(f"\nExtracting text from {file_path}:")
        result = handler.process_file(file_path)
        print(result)

if __name__ == "__main__":
    main()